# -*- coding: utf-8 -*-
"""
ANALISI CR VALUE 2.3 - FinancePlus.Tech
=======================================
Applicazione Streamlit in un unico file per:
- caricamento CR in PDF, DOCX, XLSX, CSV, TXT e immagini;
- estrazione automatica anagrafica e dati principali;
- archivio clienti SQLite con anti-duplicato per P.IVA/CF;
- più rilevazioni per lo stesso cliente;
- analisi quantitativa e score consulenziale;
- anteprima e generazione report PDF professionale;
- elenco, ricerca, download ed eliminazione report.

Avvio:
    pip install streamlit pandas openpyxl pymupdf reportlab pillow python-docx pytesseract
    streamlit run Analisi_CR_VALUE_2_1.py

Nota: per OCR installare anche Tesseract nel sistema operativo.
"""
from __future__ import annotations

import base64
import hashlib
import html
import io
import json
import re
import sqlite3
import uuid
import zipfile
from dataclasses import dataclass, asdict
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

import streamlit as st

try:
    import fitz
except Exception:
    fitz = None
try:
    import pandas as pd
except Exception:
    pd = None
try:
    from PIL import Image
except Exception:
    Image = None
try:
    import pytesseract
except Exception:
    pytesseract = None
try:
    from docx import Document
except Exception:
    Document = None
try:
    from pptx import Presentation
except Exception:
    Presentation = None

from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY, TA_LEFT
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak,
    KeepTogether, Image as RLImage
)

APP_NAME = "Analisi CR"
APP_VERSION = "2.3"
SITE = "www.financeplus.tech"
BASE_DIR = Path(__file__).resolve().parent
DATA_DIR = BASE_DIR / "analisi_cr_value_2_1_data"
CLIENTS_DIR = DATA_DIR / "clienti"
DB_PATH = DATA_DIR / "analisi_cr.sqlite3"
TEMP_DIR = DATA_DIR / "temp"

NAVY = "#071A2B"
BLUE = "#123D66"
COPPER = "#C9793D"
CYAN = "#2EA9C7"
LIGHT = "#F3F6F9"
GREEN = "#16805A"
AMBER = "#D99A22"
RED = "#C74242"
GRAY = "#657383"


def ensure_dirs() -> None:
    for p in (DATA_DIR, CLIENTS_DIR, TEMP_DIR):
        p.mkdir(parents=True, exist_ok=True)


def db() -> sqlite3.Connection:
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    return conn


def init_db() -> None:
    """Inizializza il database e migra automaticamente le versioni precedenti.

    La P.IVA/CF non e piu obbligatoria: identity_key e una chiave tecnica
    usata solo per l'anti-duplicato e non viene mostrata all'utente.
    """
    ensure_dirs()
    with db() as conn:
        existing = conn.execute("SELECT name FROM sqlite_master WHERE type='table' AND name='clients'").fetchone()
        if existing:
            cols = {r["name"]: r for r in conn.execute("PRAGMA table_info(clients)").fetchall()}
            needs_migration = "identity_key" not in cols or int(cols.get("vat_cf", {"notnull": 0})["notnull"] or 0) == 1
            if needs_migration:
                conn.execute("PRAGMA foreign_keys=OFF")
                conn.executescript("""
                ALTER TABLE clients RENAME TO clients_legacy;
                CREATE TABLE clients(
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    identity_key TEXT NOT NULL UNIQUE,
                    vat_cf TEXT DEFAULT '',
                    company_name TEXT NOT NULL DEFAULT '',
                    tax_code TEXT DEFAULT '',
                    address TEXT DEFAULT '',
                    city TEXT DEFAULT '',
                    province TEXT DEFAULT '',
                    cciaa TEXT DEFAULT '',
                    lei TEXT DEFAULT '',
                    created_at TEXT NOT NULL,
                    updated_at TEXT NOT NULL
                );
                INSERT INTO clients(id,identity_key,vat_cf,company_name,tax_code,address,city,province,cciaa,lei,created_at,updated_at)
                SELECT id,
                       CASE WHEN TRIM(COALESCE(vat_cf,''))<>'' THEN UPPER(REPLACE(vat_cf,' ','')) ELSE 'LEGACY-'||id END,
                       CASE WHEN vat_cf LIKE 'AUTO-%' OR vat_cf LIKE 'LEGACY-%' THEN '' ELSE COALESCE(vat_cf,'') END,
                       COALESCE(company_name,''),COALESCE(tax_code,''),COALESCE(address,''),COALESCE(city,''),
                       COALESCE(province,''),COALESCE(cciaa,''),COALESCE(lei,''),created_at,updated_at
                FROM clients_legacy;
                DROP TABLE clients_legacy;
                """)
                conn.execute("PRAGMA foreign_keys=ON")
        conn.executescript("""
        CREATE TABLE IF NOT EXISTS clients(
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            identity_key TEXT NOT NULL UNIQUE,
            vat_cf TEXT DEFAULT '',
            company_name TEXT NOT NULL DEFAULT '',
            tax_code TEXT DEFAULT '',
            address TEXT DEFAULT '',
            city TEXT DEFAULT '',
            province TEXT DEFAULT '',
            cciaa TEXT DEFAULT '',
            lei TEXT DEFAULT '',
            created_at TEXT NOT NULL,
            updated_at TEXT NOT NULL
        );
        CREATE TABLE IF NOT EXISTS reports(
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            client_id INTEGER NOT NULL,
            report_date TEXT NOT NULL,
            source_name TEXT NOT NULL,
            source_hash TEXT NOT NULL,
            source_path TEXT NOT NULL,
            report_path TEXT NOT NULL,
            period TEXT,
            score REAL,
            rating TEXT,
            analysis_json TEXT,
            created_at TEXT NOT NULL,
            FOREIGN KEY(client_id) REFERENCES clients(id)
        );
        CREATE INDEX IF NOT EXISTS idx_reports_client ON reports(client_id);
        CREATE INDEX IF NOT EXISTS idx_reports_hash ON reports(source_hash);
        CREATE INDEX IF NOT EXISTS idx_clients_identity ON clients(identity_key);
        """)


def now_iso() -> str:
    return datetime.now().isoformat(timespec="seconds")


def it_date(value: Optional[str] = None, with_time: bool = False) -> str:
    if value:
        try:
            dt = datetime.fromisoformat(value)
        except Exception:
            return value
    else:
        dt = datetime.now()
    return dt.strftime("%d-%m-%Y %H:%M" if with_time else "%d-%m-%Y")


def safe_name(text: str) -> str:
    text = re.sub(r"[^A-Za-z0-9._-]+", "_", (text or "CLIENTE").strip())
    return re.sub(r"_+", "_", text).strip("_") or "CLIENTE"


def sha256(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def money(v: Any) -> str:
    try:
        x = float(v or 0)
    except Exception:
        x = 0
    s = f"{x:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")
    return f"€ {s}"


def pct(v: Any) -> str:
    try:
        x = float(v or 0)
    except Exception:
        x = 0
    return f"{x:.1f}%".replace(".", ",")


@dataclass
class ClientData:
    company_name: str = ""
    vat_cf: str = ""
    tax_code: str = ""
    address: str = ""
    city: str = ""
    province: str = ""
    cciaa: str = ""
    lei: str = ""
    period: str = ""


@dataclass
class CRAnalysis:
    metrics: Dict[str, float]
    banks: List[Dict[str, Any]]
    categories: List[Dict[str, Any]]
    monthly: List[Dict[str, Any]]
    anomalies: List[str]
    recommendations: List[str]
    score: float
    rating: str
    narrative: str
    warnings: List[str]


class DocumentExtractor:
    """Estrae testo dai formati noti e accetta comunque qualsiasi file.

    I formati non interpretabili vengono archiviati senza bloccare il flusso.
    """
    def extract(self, name: str, data: bytes) -> Tuple[str, List[str]]:
        ext = Path(name).suffix.lower()
        warnings: List[str] = []
        if ext == ".pdf":
            return self._pdf(data)
        if ext == ".docx":
            if Document is None:
                return "", ["python-docx non installato: il file sara comunque archiviato."]
            try:
                d = Document(io.BytesIO(data))
                parts = [p.text for p in d.paragraphs]
                for t in d.tables:
                    for r in t.rows:
                        parts.append("\t".join(c.text for c in r.cells))
                return "\n".join(parts), warnings
            except Exception as e:
                return "", [f"DOCX non interpretabile ({e}); il file sara comunque archiviato."]
        if ext == ".pptx":
            if Presentation is None:
                return "", ["python-pptx non installato: il file sara comunque archiviato."]
            try:
                prs = Presentation(io.BytesIO(data))
                parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            parts.append(shape.text)
                return "\n".join(parts), warnings
            except Exception as e:
                return "", [f"PPTX non interpretabile ({e}); il file sara comunque archiviato."]
        if ext in {".xlsx", ".xls", ".xlsm", ".xlsb", ".ods", ".csv", ".tsv"}:
            if pd is None:
                return "", ["pandas non installato: il file sara comunque archiviato."]
            try:
                if ext in {".csv", ".tsv"}:
                    frame = pd.read_csv(io.BytesIO(data), sep=("\t" if ext == ".tsv" else None), engine="python")
                else:
                    sheets = pd.read_excel(io.BytesIO(data), sheet_name=None)
                    return "\n\n".join(f"FOGLIO: {k}\n{v.to_csv(index=False)}" for k, v in sheets.items()), warnings
                return frame.to_csv(index=False), warnings
            except Exception as e:
                return "", [f"File tabellare non interpretabile ({e}); il file sara comunque archiviato."]
        if ext in {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".webp", ".gif"}:
            if Image is None or pytesseract is None:
                return "", ["OCR non disponibile: il file immagine sara comunque archiviato."]
            try:
                img = Image.open(io.BytesIO(data)).convert("RGB")
                return pytesseract.image_to_string(img, lang="ita+eng"), warnings
            except Exception as e:
                return "", [f"OCR non riuscito ({e}); il file sara comunque archiviato."]
        if ext in {".rtf", ".txt", ".log", ".md", ".json", ".xml", ".html", ".htm", ".yaml", ".yml", ".ini", ".cfg"}:
            text = data.decode("utf-8", errors="replace")
            if ext == ".rtf":
                text = re.sub(r"\\[a-z]+-?\d* ?|[{}]", " ", text, flags=re.I)
            if ext in {".html", ".htm", ".xml"}:
                text = re.sub(r"<[^>]+>", " ", text)
            return html.unescape(text), warnings
        if ext == ".odt":
            try:
                with zipfile.ZipFile(io.BytesIO(data)) as zf:
                    raw = zf.read("content.xml").decode("utf-8", errors="replace")
                return html.unescape(re.sub(r"<[^>]+>", " ", raw)), warnings
            except Exception as e:
                return "", [f"ODT non interpretabile ({e}); il file sara comunque archiviato."]
        # Vecchi DOC, archivi, email, CAD e ogni altro formato: upload consentito.
        try:
            sample = data.decode("utf-8")
            printable = sum(ch.isprintable() or ch in "\n\r\t" for ch in sample) / max(len(sample), 1)
            if printable > 0.85:
                return sample, warnings
        except Exception:
            pass
        return "", [f"Formato {ext or 'senza estensione'} accettato e archiviabile, ma non disponibile per l'estrazione automatica del testo."]

    def _pdf(self, data: bytes) -> Tuple[str, List[str]]:
        if fitz is None:
            return "", ["PyMuPDF non installato."]
        warnings: List[str] = []
        chunks: List[str] = []
        doc = fitz.open(stream=data, filetype="pdf")
        for i, page in enumerate(doc):
            txt = page.get_text("text") or ""
            if len(txt.strip()) < 30 and Image is not None and pytesseract is not None:
                try:
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                    img = Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGB")
                    txt = pytesseract.image_to_string(img, lang="ita+eng")
                except Exception as e:
                    warnings.append(f"OCR pagina {i+1} non riuscito: {e}")
            chunks.append(txt)
        doc.close()
        if not any(c.strip() for c in chunks):
            warnings.append("Nessun testo estratto: verificare OCR e qualità del documento.")
        return "\n\f\n".join(chunks), warnings


class CRParser:
    BANK_WORDS = ("BANCA", "BANCO", "INTESA", "UNICREDIT", "BPER", "BPM", "BCC", "IFIS", "CREDITO", "MEDIOCREDITO", "SANTANDER", "BNL")

    def parse(self, text: str, warnings: List[str]) -> Tuple[ClientData, CRAnalysis]:
        client = self._client(text)
        amounts = self._amounts(text)
        categories = self._categories(text)
        banks = self._banks(text)
        monthly = self._monthly(text)

        accordato = amounts.get("accordato", 0.0)
        utilizzato = amounts.get("utilizzato", 0.0)
        sconfini = amounts.get("sconfini", 0.0)
        sofferenze = amounts.get("sofferenze", 0.0)
        scaduti = amounts.get("scaduti", 0.0)
        garanzie = amounts.get("garanzie", 0.0)
        usage = (utilizzato / accordato * 100) if accordato > 0 else 0

        anomalies: List[str] = []
        deductions = 0.0
        if sofferenze > 0:
            anomalies.append("Presenza di sofferenze o crediti deteriorati.")
            deductions += 45
        if scaduti > 0:
            anomalies.append("Presenza di scaduti/past due.")
            deductions += min(25, 8 + scaduti / max(utilizzato, 1) * 100)
        if sconfini > 0:
            anomalies.append("Sconfinamenti rilevati nel periodo.")
            deductions += min(20, 5 + sconfini / max(accordato, 1) * 250)
        if usage >= 95:
            anomalies.append("Saturazione degli affidamenti molto elevata.")
            deductions += 12
        elif usage >= 85:
            anomalies.append("Saturazione degli affidamenti elevata.")
            deductions += 6
        if len(banks) >= 8:
            anomalies.append("Elevata frammentazione degli intermediari segnalanti.")
            deductions += 4
        if not anomalies:
            anomalies.append("Non emergono anomalie quantitative rilevanti dai dati riconosciuti automaticamente.")

        score = max(0.0, min(100.0, 100.0 - deductions))
        if score >= 90:
            rating = "OTTIMO"
        elif score >= 80:
            rating = "BUONO"
        elif score >= 65:
            rating = "ADEGUATO"
        elif score >= 50:
            rating = "VULNERABILE"
        else:
            rating = "CRITICO"

        recs = []
        if usage >= 85:
            recs.append("Ridurre la saturazione media delle linee, preservando un margine disponibile almeno del 15-20%.")
        if sconfini > 0:
            recs.append("Azzerare gli sconfinamenti e verificare eventuali disallineamenti di valuta o plafond.")
        if scaduti > 0:
            recs.append("Analizzare e regolarizzare le esposizioni scadute, documentando eventuali contestazioni.")
        if len(banks) > 6:
            recs.append("Valutare il consolidamento delle relazioni bancarie e la razionalizzazione delle linee minori.")
        recs += [
            "Confrontare la CR con bilanci, situazione contabile aggiornata ed estratti conto.",
            "Verificare manualmente ogni importo prima dell'utilizzo professionale del report.",
        ]
        narrative = (
            f"L'analisi automatica rileva affidamenti per {money(accordato)}, utilizzi per {money(utilizzato)} "
            f"e un tasso di utilizzo del {pct(usage)}. Il punteggio consulenziale è {score:.0f}/100, "
            f"con giudizio {rating}."
        )
        metrics = {
            "accordato": accordato, "utilizzato": utilizzato, "sconfini": sconfini,
            "sofferenze": sofferenze, "scaduti": scaduti, "garanzie": garanzie,
            "utilizzo_pct": usage, "banche": float(len(banks)), "mesi": float(len(monthly) or 1),
        }
        return client, CRAnalysis(metrics, banks, categories, monthly, anomalies, recs, score, rating, narrative, warnings)

    def _client(self, text: str) -> ClientData:
        t = re.sub(r"[ \t]+", " ", text)
        lines = [x.strip() for x in t.splitlines() if x.strip()]
        company = ""
        company_patterns = [
            r"(?:denominazione|soggetto della visura|ragione sociale)\s*[:\-]?\s*([A-Z0-9À-Ù' .&-]{4,100})",
            r"\b([A-Z][A-Z0-9À-Ù' .&-]{3,80}\s(?:S\.?R\.?L\.?|S\.?P\.?A\.?|S\.?A\.?S\.?|S\.?N\.?C\.?))\b",
        ]
        for p in company_patterns:
            m = re.search(p, t, re.I)
            if m:
                company = m.group(1).strip(" -:")
                break
        if not company:
            for line in lines[:40]:
                if any(x in line.upper() for x in ("S.R.L", "SRL", "S.P.A", "SPA")):
                    company = line[:100]
                    break
        vat = ""
        for p in [r"(?:partita iva|p\.?iva|codice fiscale\s*/\s*p\.?iva)\s*[:\-]?\s*([A-Z0-9]{11,16})", r"\b([0-9]{11})\b"]:
            m = re.search(p, t, re.I)
            if m:
                vat = m.group(1)
                break
        tax = ""
        m = re.search(r"codice fiscale\s*[:\-]?\s*([A-Z0-9]{11,16})", t, re.I)
        if m: tax = m.group(1)
        city = ""
        province = ""
        m = re.search(r"(?:comune|sede)\s*[:\-]?\s*([A-ZÀ-Ù' ]{2,50})(?:\s+\(?([A-Z]{2})\)?)?", t, re.I)
        if m:
            city = m.group(1).strip()
            province = (m.group(2) or "").upper()
        address = ""
        m = re.search(r"(?:indirizzo|via|sede legale)\s*[:\-]?\s*((?:VIA|VIALE|PIAZZA|CORSO|LOC\.)?[^\n]{5,100})", t, re.I)
        if m: address = m.group(1).strip()
        cciaa = ""
        m = re.search(r"(?:cciaa|rea)\s*[:\-]?\s*([A-Z0-9-]{4,20})", t, re.I)
        if m: cciaa = m.group(1)
        lei = ""
        m = re.search(r"(?:codice lei|lei)\s*[:\-]?\s*([A-Z0-9]{18,20})", t, re.I)
        if m: lei = m.group(1)
        period = ""
        m = re.search(r"(?:periodo censito|periodo di riferimento)\s*[:\-]?\s*([^\n]{3,40})", t, re.I)
        if m: period = m.group(1).strip()
        return ClientData(company or "CLIENTE DA VERIFICARE", vat or tax, tax, address, city, province, cciaa, lei, period)

    def _num(self, raw: str) -> float:
        s = raw.replace("€", "").replace(" ", "").strip()
        if "." in s and "," in s:
            s = s.replace(".", "").replace(",", ".")
        elif "," in s:
            s = s.replace(".", "").replace(",", ".")
        else:
            s = s.replace(".", "")
        try: return float(s)
        except Exception: return 0.0

    def _find_amount(self, text: str, labels: List[str]) -> float:
        best = 0.0
        for label in labels:
            pat = rf"{label}[^\d€]{{0,50}}€?\s*([0-9][0-9. ]*(?:,[0-9]{{1,2}})?)"
            for m in re.finditer(pat, text, re.I):
                best = max(best, self._num(m.group(1)))
        return best

    def _amounts(self, text: str) -> Dict[str, float]:
        return {
            "accordato": self._find_amount(text, ["accordato operativo", "attuali affidamenti", "monte affidamenti", "accordato"]),
            "utilizzato": self._find_amount(text, ["attuali utilizzi", "monte utilizzi", "utilizzato"]),
            "sconfini": self._find_amount(text, ["sconfini / insoluti", "sconfinamento", "sconfini"]),
            "sofferenze": self._find_amount(text, ["sofferenze", "crediti passati a perdita"]),
            "scaduti": self._find_amount(text, ["scaduti", "past due", "scaduto/sconfinato"]),
            "garanzie": self._find_amount(text, ["valore complessivo delle garanzie", "garanzie"]),
        }

    def _categories(self, text: str) -> List[Dict[str, Any]]:
        out = []
        mapping = {
            "Rischi autoliquidanti": "autoliquidanti",
            "Rischi a scadenza": "a scadenza",
            "Rischi a revoca": "a revoca",
            "Crediti di firma": "crediti di firma",
        }
        for name, label in mapping.items():
            amount = self._find_amount(text, [label])
            if amount > 0:
                out.append({"categoria": name, "importo": amount})
        return out

    def _banks(self, text: str) -> List[Dict[str, Any]]:
        rows: Dict[str, float] = {}
        for line in text.splitlines():
            u = line.upper().strip()
            if len(u) < 6 or not any(w in u for w in self.BANK_WORDS):
                continue
            name_match = re.search(r"([A-ZÀ-Ù0-9' .&-]{5,100}?(?:S\.P\.A\.|SPA|SOCIETA' COOPERATIVA|SOCIETÀ COOPERATIVA|BANCA|BANCO))", u)
            nums = re.findall(r"€?\s*([0-9][0-9. ]*(?:,[0-9]{1,2})?)", line)
            if name_match:
                name = re.sub(r"\s+", " ", name_match.group(1)).strip(" -")
                vals = [self._num(n) for n in nums]
                amount = max(vals) if vals else 0.0
                if name and name not in {"BANCA", "BANCO"}:
                    rows[name] = max(rows.get(name, 0), amount)
        return [{"banca": k, "importo": v} for k, v in sorted(rows.items(), key=lambda x: x[1], reverse=True)[:20]]

    def _monthly(self, text: str) -> List[Dict[str, Any]]:
        months = {"gen", "feb", "mar", "apr", "mag", "giu", "lug", "ago", "set", "ott", "nov", "dic"}
        out = []
        seen = set()
        for line in text.splitlines():
            m = re.search(r"\b(gen|feb|mar|apr|mag|giu|lug|ago|set|ott|nov|dic)\s+(20\d{2})\b", line, re.I)
            if m:
                key = f"{m.group(1).lower()} {m.group(2)}"
                if key not in seen:
                    seen.add(key)
                    out.append({"periodo": key})
        return out[-36:]


def _normalize_identity(value: str) -> str:
    return re.sub(r"[^A-Z0-9]", "", (value or "").upper())


def _identity_key(c: ClientData) -> str:
    fiscal = _normalize_identity(c.vat_cf or c.tax_code)
    if fiscal:
        return "FISCAL-" + fiscal
    company = re.sub(r"[^A-Z0-9]", "", (c.company_name or "").upper())
    city = re.sub(r"[^A-Z0-9]", "", (c.city or "").upper())
    if company:
        return "NAME-" + company + ("-" + city if city else "")
    return "ANON-" + uuid.uuid4().hex


def upsert_client(c: ClientData) -> Tuple[int, bool]:
    # Nessun alert bloccante su P.IVA o CF. Basta una denominazione, anche generica.
    if not (c.company_name or "").strip():
        c.company_name = "Cliente senza denominazione"
    identity = _identity_key(c)
    vat = (c.vat_cf or "").strip()
    tax = (c.tax_code or "").strip()
    with db() as conn:
        row = conn.execute("SELECT id FROM clients WHERE identity_key=?", (identity,)).fetchone()
        if row:
            conn.execute("""UPDATE clients SET vat_cf=?,company_name=?,tax_code=?,address=?,city=?,province=?,cciaa=?,lei=?,updated_at=? WHERE id=?""",
                         (vat, c.company_name, tax, c.address, c.city, c.province, c.cciaa, c.lei, now_iso(), row["id"]))
            return int(row["id"]), False
        cur = conn.execute("""INSERT INTO clients(identity_key,vat_cf,company_name,tax_code,address,city,province,cciaa,lei,created_at,updated_at)
                              VALUES(?,?,?,?,?,?,?,?,?,?,?)""",
                           (identity, vat, c.company_name, tax, c.address, c.city, c.province, c.cciaa, c.lei, now_iso(), now_iso()))
        return int(cur.lastrowid), True


def client_folder(c: ClientData) -> Path:
    fiscal = safe_name(c.vat_cf or c.tax_code)
    suffix = fiscal if fiscal and fiscal != "CLIENTE" else safe_name(c.city or "SENZA_CF")
    p = CLIENTS_DIR / f"{safe_name(c.company_name)}_{suffix}"
    (p / "CR_originali").mkdir(parents=True, exist_ok=True)
    (p / "Report_elaborati").mkdir(parents=True, exist_ok=True)
    return p


def save_report(c: ClientData, a: CRAnalysis, source_name: str, source: bytes, pdf: bytes) -> Tuple[bool, str]:
    client_id, created = upsert_client(c)
    h = sha256(source)
    with db() as conn:
        dup = conn.execute("SELECT id,report_date FROM reports WHERE client_id=? AND source_hash=?", (client_id, h)).fetchone()
    if dup:
        raise ValueError(f"Documento già archiviato il {dup['report_date']}.")
    folder = client_folder(c)
    stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    ext = Path(source_name).suffix or ".bin"
    src_path = folder / "CR_originali" / f"{safe_name(c.company_name)}_CR_ORIGINALE_{stamp}{ext}"
    pdf_path = folder / "Report_elaborati" / f"{safe_name(c.company_name)}_ANALISI_CR_{stamp}.pdf"
    src_path.write_bytes(source)
    pdf_path.write_bytes(pdf)
    with db() as conn:
        conn.execute("""INSERT INTO reports(client_id,report_date,source_name,source_hash,source_path,report_path,period,score,rating,analysis_json,created_at)
                        VALUES(?,?,?,?,?,?,?,?,?,?,?)""",
                     (client_id, it_date(), source_name, h, str(src_path), str(pdf_path), c.period,
                      a.score, a.rating, json.dumps(asdict(a), ensure_ascii=False), now_iso()))
    return created, str(pdf_path)


def report_rows(search: str = ""):
    with db() as conn:
        q = """SELECT r.*,c.company_name,c.vat_cf,c.city,c.province FROM reports r JOIN clients c ON c.id=r.client_id"""
        if search.strip():
            x = f"%{search.strip()}%"
            return conn.execute(q + " WHERE c.company_name LIKE ? OR c.vat_cf LIKE ? OR r.rating LIKE ? ORDER BY r.id DESC", (x, x, x)).fetchall()
        return conn.execute(q + " ORDER BY r.id DESC").fetchall()


def delete_report(report_id: int) -> bool:
    with db() as conn:
        row = conn.execute("SELECT * FROM reports WHERE id=?", (report_id,)).fetchone()
        if not row: return False
        for k in ("source_path", "report_path"):
            try: Path(row[k]).unlink(missing_ok=True)
            except Exception: pass
        conn.execute("DELETE FROM reports WHERE id=?", (report_id,))
    return True


def build_report(c: ClientData, a: CRAnalysis) -> bytes:
    buf = io.BytesIO()
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name="CoverTitle", parent=styles["Title"], fontName="Helvetica-Bold", fontSize=27, leading=32, textColor=colors.HexColor(NAVY), alignment=TA_CENTER, spaceAfter=10))
    styles.add(ParagraphStyle(name="H1x", parent=styles["Heading1"], fontName="Helvetica-Bold", fontSize=17, leading=21, textColor=colors.HexColor(NAVY), spaceAfter=8))
    styles.add(ParagraphStyle(name="H2x", parent=styles["Heading2"], fontName="Helvetica-Bold", fontSize=11, leading=14, textColor=colors.HexColor(COPPER), spaceAfter=5))
    styles.add(ParagraphStyle(name="Bodyx", parent=styles["BodyText"], fontName="Helvetica", fontSize=8.5, leading=12, alignment=TA_JUSTIFY, textColor=colors.HexColor("#263544")))
    styles.add(ParagraphStyle(name="Smallx", parent=styles["BodyText"], fontSize=7, leading=9, textColor=colors.HexColor(GRAY)))
    styles.add(ParagraphStyle(name="Scorex", parent=styles["Title"], fontName="Helvetica-Bold", fontSize=31, leading=33, textColor=colors.HexColor(COPPER), alignment=TA_CENTER))

    def header_footer(canvas, doc):
        canvas.saveState()
        w, h = A4
        canvas.setFillColor(colors.HexColor(NAVY))
        canvas.rect(0, h-13*mm, w, 13*mm, fill=1, stroke=0)
        canvas.setFillColor(colors.white)
        canvas.setFont("Helvetica-Bold", 9)
        canvas.drawString(15*mm, h-8.5*mm, "FINANCEPLUS.TECH | ANALISI CR VALUE 2.3")
        canvas.setStrokeColor(colors.HexColor(COPPER))
        canvas.line(15*mm, 15*mm, w-15*mm, 15*mm)
        canvas.setFillColor(colors.HexColor(GRAY))
        canvas.setFont("Helvetica", 7)
        canvas.drawString(15*mm, 9*mm, SITE)
        canvas.drawRightString(w-15*mm, 9*mm, f"Pagina {doc.page}")
        canvas.restoreState()

    doc = SimpleDocTemplate(buf, pagesize=A4, rightMargin=16*mm, leftMargin=16*mm, topMargin=20*mm, bottomMargin=20*mm,
                            title=f"Analisi CR - {c.company_name}", author="FinancePlus.Tech")
    story = []

    story += [Spacer(1, 26*mm), Paragraph("ANALISI CR AVANZATA", styles["CoverTitle"]), Spacer(1, 8*mm),
              Paragraph(html.escape(c.company_name), styles["H1x"]), Spacer(1, 4*mm)]
    cover = [
        ["Partita IVA / CF", c.vat_cf or c.tax_code or "N/D"], ["Comune", c.city or "N/D"],
        ["Provincia", c.province or "N/D"], ["Periodo censito", c.period or "N/D"],
        ["Data elaborazione", it_date(with_time=True)],
    ]
    story += [Table(cover, colWidths=[48*mm, 105*mm], style=TableStyle([
        ("BACKGROUND", (0,0),(0,-1), colors.HexColor(LIGHT)), ("TEXTCOLOR",(0,0),(0,-1),colors.HexColor(NAVY)),
        ("FONTNAME",(0,0),(0,-1),"Helvetica-Bold"),("FONTNAME",(1,0),(1,-1),"Helvetica"),
        ("FONTSIZE",(0,0),(-1,-1),9),("GRID",(0,0),(-1,-1),0.3,colors.HexColor("#C9D3DC")),
        ("VALIGN",(0,0),(-1,-1),"MIDDLE"),("TOPPADDING",(0,0),(-1,-1),7),("BOTTOMPADDING",(0,0),(-1,-1),7)
    ])), Spacer(1, 18*mm), Paragraph("Report consulenziale interno. Verificare i dati con la Centrale Rischi originale prima di ogni utilizzo.", styles["Smallx"]), PageBreak()]

    story += [Paragraph("1. Executive summary", styles["H1x"]), Paragraph(a.narrative, styles["Bodyx"]), Spacer(1, 7*mm)]
    cards = [[Paragraph("CR SCORE", styles["Smallx"]), Paragraph(f"{a.score:.0f}/100", styles["Scorex"]), Paragraph(a.rating, styles["H2x"])],
             [Paragraph("AFFIDAMENTI", styles["Smallx"]), Paragraph(money(a.metrics['accordato']), styles["H1x"]), Paragraph("Accordato rilevato", styles["Smallx"])],
             [Paragraph("UTILIZZI", styles["Smallx"]), Paragraph(money(a.metrics['utilizzato']), styles["H1x"]), Paragraph(pct(a.metrics['utilizzo_pct']) + " saturazione", styles["Smallx"])]]
    story += [Table([cards], colWidths=[57*mm]*3, style=TableStyle([
        ("BACKGROUND",(0,0),(-1,-1),colors.HexColor(LIGHT)),("BOX",(0,0),(-1,-1),0.6,colors.HexColor("#CCD5DE")),
        ("INNERGRID",(0,0),(-1,-1),0.4,colors.white),("VALIGN",(0,0),(-1,-1),"MIDDLE"),("ALIGN",(0,0),(-1,-1),"CENTER"),
        ("TOPPADDING",(0,0),(-1,-1),8),("BOTTOMPADDING",(0,0),(-1,-1),8)
    ])), Spacer(1, 8*mm), Paragraph("Elementi di attenzione", styles["H2x"])]
    for x in a.anomalies: story.append(Paragraph("• " + html.escape(x), styles["Bodyx"]))
    story.append(PageBreak())

    story += [Paragraph("2. Anagrafica del soggetto", styles["H1x"])]
    anag = [["Campo", "Dato"], ["Denominazione", c.company_name], ["P.IVA / CF", c.vat_cf], ["Codice fiscale", c.tax_code],
            ["Indirizzo", c.address], ["Comune", c.city], ["Provincia", c.province], ["CCIAA / REA", c.cciaa], ["Codice LEI", c.lei], ["Periodo", c.period]]
    story += [_table(anag, [48*mm, 120*mm]), Spacer(1, 8*mm), Paragraph("I campi restano modificabili nell'applicazione prima della generazione definitiva.", styles["Bodyx"]), PageBreak()]

    story += [Paragraph("3. Quadro quantitativo", styles["H1x"])]
    q = [["Indicatore", "Valore", "Lettura"], ["Accordato", money(a.metrics['accordato']), "Monte affidamenti riconosciuto"],
         ["Utilizzato", money(a.metrics['utilizzato']), "Esposizione utilizzata"], ["Saturazione", pct(a.metrics['utilizzo_pct']), "Utilizzato / accordato"],
         ["Sconfini", money(a.metrics['sconfini']), "Importi oltre limite"], ["Scaduti", money(a.metrics['scaduti']), "Past due riconosciuto"],
         ["Sofferenze", money(a.metrics['sofferenze']), "Deteriorato riconosciuto"], ["Garanzie", money(a.metrics['garanzie']), "Valore garanzie"]]
    story += [_table(q, [52*mm, 42*mm, 74*mm]), Spacer(1, 7*mm), Paragraph("L'assenza di un importo può dipendere dalla struttura del documento e non equivale automaticamente ad assenza della voce.", styles["Smallx"]), PageBreak()]

    story += [Paragraph("4. Forme tecniche", styles["H1x"])]
    cat = [["Categoria", "Importo riconosciuto"]] + [[x['categoria'], money(x['importo'])] for x in a.categories]
    if len(cat) == 1: cat.append(["Nessuna forma tecnica riclassificata automaticamente", "N/D"])
    story += [_table(cat, [105*mm, 63*mm]), Spacer(1, 8*mm), Paragraph("La lettura per forma tecnica deve essere confrontata con accordato operativo, utilizzato, scaduto e sconfinamento per ciascun intermediario.", styles["Bodyx"]), PageBreak()]

    story += [Paragraph("5. Distinta intermediari", styles["H1x"])]
    bank = [["Intermediario riconosciuto", "Importo massimo associato"]] + [[x['banca'], money(x['importo'])] for x in a.banks]
    if len(bank) == 1: bank.append(["Nessun intermediario riconosciuto automaticamente", "N/D"])
    story += [_table(bank, [120*mm, 48*mm]), PageBreak()]

    story += [Paragraph("6. Evoluzione temporale", styles["H1x"]), Paragraph(f"Mesi/periodi riconosciuti: {int(a.metrics['mesi'])}.", styles["Bodyx"]), Spacer(1, 5*mm)]
    periods = [["#", "Periodo"]] + [[str(i+1), x['periodo']] for i, x in enumerate(a.monthly)]
    if len(periods) == 1: periods.append(["-", "Serie storica non riclassificata automaticamente"])
    story += [_table(periods, [20*mm, 148*mm]), PageBreak()]

    story += [Paragraph("7. Eventi negativi e alert", styles["H1x"])]
    alert_rows = [["Fattore", "Esito"], ["Sofferenze", "PRESENTI" if a.metrics['sofferenze'] > 0 else "Non rilevate"],
                  ["Scaduti/past due", "PRESENTI" if a.metrics['scaduti'] > 0 else "Non rilevati"],
                  ["Sconfini", "PRESENTI" if a.metrics['sconfini'] > 0 else "Non rilevati"],
                  ["Saturazione >= 95%", "SI" if a.metrics['utilizzo_pct'] >= 95 else "NO"]]
    story += [_table(alert_rows, [88*mm, 80*mm]), Spacer(1, 8*mm)]
    for x in a.anomalies: story.append(Paragraph("• " + html.escape(x), styles["Bodyx"]))
    story.append(PageBreak())

    story += [Paragraph("8. Valutazione consulenziale", styles["H1x"]), Paragraph(f"Punteggio: {a.score:.0f}/100 - Giudizio: {a.rating}", styles["H2x"]),
              Paragraph(a.narrative, styles["Bodyx"]), Spacer(1, 8*mm), Paragraph("Scala interpretativa", styles["H2x"])]
    scale = [["Fascia", "Giudizio", "Interpretazione"], ["90-100", "Ottimo", "Profilo andamentale robusto"], ["80-89", "Buono", "Profilo positivo con aspetti da monitorare"],
             ["65-79", "Adeguato", "Finanziabilità da integrare con bilancio e flussi"], ["50-64", "Vulnerabile", "Presenza di criticità da correggere"], ["0-49", "Critico", "Rischio elevato e necessità di approfondimento"]]
    story += [_table(scale, [28*mm, 35*mm, 105*mm]), PageBreak()]

    story += [Paragraph("9. Azioni consigliate", styles["H1x"])]
    for i, x in enumerate(a.recommendations, 1):
        story.append(KeepTogether([Paragraph(f"{i}. {html.escape(x)}", styles["Bodyx"]), Spacer(1, 3*mm)]))
    story.append(PageBreak())

    story += [Paragraph("10. Checklist di verifica", styles["H1x"])]
    checklist = [["Verifica", "Stato"], ["Anagrafica e P.IVA confermate", "□"], ["Periodo CR confermato", "□"], ["Accordato e utilizzato verificati", "□"],
                 ["Sconfini e scaduti verificati", "□"], ["Sofferenze/perdite verificate", "□"], ["Garanzie e coobbligazioni verificate", "□"],
                 ["Confronto con bilanci e situazione aggiornata", "□"], ["Confronto con estratti conto", "□"]]
    story += [_table(checklist, [142*mm, 26*mm]), PageBreak()]

    story += [Paragraph("11. Avvertenze metodologiche", styles["H1x"]),
              Paragraph("Il report deriva da estrazione automatica e regole euristiche. Non costituisce rating ufficiale di Banca d'Italia, MCC, banca o credit bureau, non rappresenta una delibera e non garantisce la concessione di credito.", styles["Bodyx"]), Spacer(1, 5*mm),
              Paragraph("La qualità dipende da completezza, leggibilità e struttura del documento. OCR e formati non standard possono generare omissioni o associazioni errate. Ogni dato deve essere verificato con la fonte originale.", styles["Bodyx"]), Spacer(1, 5*mm),
              Paragraph("La valutazione va integrata con bilanci, situazione contabile, flussi bancari, pregiudizievoli, assetto societario, settore, governance e finalità dell'operazione.", styles["Bodyx"]), PageBreak()]

    story += [Paragraph("12. Conclusioni", styles["H1x"]), Paragraph(a.narrative, styles["Bodyx"]), Spacer(1, 8*mm),
              Paragraph("Esito operativo", styles["H2x"]), Paragraph("Il dossier è idoneo come base di lavoro interna. Prima della consegna a banca o intermediario occorre completare la verifica documentale e motivare le eventuali anomalie.", styles["Bodyx"])]

    doc.build(story, onFirstPage=header_footer, onLaterPages=header_footer)
    return buf.getvalue()


def _table(data: List[List[Any]], widths: List[float]) -> Table:
    normalized = [[Paragraph(html.escape(str(c)) if c is not None else "", ParagraphStyle(name=f"t{id(data)}", fontName="Helvetica", fontSize=7.2, leading=9)) for c in row] for row in data]
    t = Table(normalized, colWidths=widths, repeatRows=1)
    t.setStyle(TableStyle([
        ("BACKGROUND",(0,0),(-1,0),colors.HexColor(NAVY)),("TEXTCOLOR",(0,0),(-1,0),colors.white),
        ("FONTNAME",(0,0),(-1,0),"Helvetica-Bold"),("GRID",(0,0),(-1,-1),0.35,colors.HexColor("#C9D3DC")),
        ("ROWBACKGROUNDS",(0,1),(-1,-1),[colors.white,colors.HexColor(LIGHT)]),("VALIGN",(0,0),(-1,-1),"MIDDLE"),
        ("TOPPADDING",(0,0),(-1,-1),5),("BOTTOMPADDING",(0,0),(-1,-1),5),("LEFTPADDING",(0,0),(-1,-1),4),("RIGHTPADDING",(0,0),(-1,-1),4)
    ]))
    return t


def pdf_preview(data: bytes, height: int = 850) -> None:
    encoded = base64.b64encode(data).decode("ascii")
    st.components.v1.html(f'<iframe src="data:application/pdf;base64,{encoded}" width="100%" height="{height}" style="border:1px solid #cbd5df;border-radius:10px"></iframe>', height=height+20)


def css() -> None:
    st.markdown(f"""
    <style>
    .stApp {{background:linear-gradient(180deg,#f8fafc,#eef3f7);}}
    [data-testid="stSidebar"] {{background:{NAVY};}}
    [data-testid="stSidebar"] * {{color:#f5f7fa;}}
    .fp-title {{font-size:2rem;font-weight:900;color:{NAVY};letter-spacing:-.03em;}}
    .fp-sub {{color:{GRAY};margin-top:-8px;margin-bottom:18px;}}
    .fp-card {{background:white;border:1px solid #dce4eb;border-radius:14px;padding:18px;box-shadow:0 4px 18px rgba(7,26,43,.06);}}
    .fp-brand {{font-size:24px;font-weight:900;line-height:1.1}} .fp-brand span {{color:{COPPER};}}
    .stButton>button {{border-radius:9px;font-weight:700;}}
    </style>""", unsafe_allow_html=True)


def header(title: str, subtitle: str) -> None:
    st.markdown(f'<div class="fp-title">{html.escape(title)}</div><div class="fp-sub">{html.escape(subtitle)}</div>', unsafe_allow_html=True)



def request_navigation(page: str) -> None:
    """Accoda il cambio pagina senza modificare il widget dopo la sua creazione."""
    st.session_state["_requested_nav_page"] = page


def apply_pending_navigation() -> None:
    """Applica il cambio pagina prima di istanziare il radio di navigazione."""
    requested = st.session_state.pop("_requested_nav_page", None)
    if requested in {"Dashboard", "Inserisci CR", "Elenco Report"}:
        st.session_state["nav_page"] = requested

def page_dashboard() -> None:
    header("Dashboard", "Archivio Centrale Rischi e report FinancePlus.Tech")
    with db() as conn:
        clients = conn.execute("SELECT COUNT(*) n FROM clients").fetchone()["n"]
        reports = conn.execute("SELECT COUNT(*) n FROM reports").fetchone()["n"]
        avg = conn.execute("SELECT AVG(score) n FROM reports").fetchone()["n"] or 0
        recent = conn.execute("SELECT COUNT(*) n FROM reports WHERE created_at >= datetime('now','-30 day')").fetchone()["n"]
    c1,c2,c3,c4 = st.columns(4)
    c1.metric("Clienti", clients); c2.metric("Report", reports); c3.metric("Score medio", f"{avg:.0f}/100"); c4.metric("Ultimi 30 giorni", recent)
    st.markdown("### Flusso operativo")
    st.info("Carica un file di qualsiasi formato → estrazione automatica quando disponibile → verifica anagrafica → analisi → anteprima PDF → salvataggio nell'archivio cliente.")
    q1, q2 = st.columns(2)
    q1.button(
        "Inserisci nuovo documento",
        type="primary",
        use_container_width=True,
        on_click=request_navigation,
        args=("Inserisci CR",),
    )
    q2.button(
        "Apri elenco report",
        use_container_width=True,
        on_click=request_navigation,
        args=("Elenco Report",),
    )
    rows = report_rows()[:8]
    if rows:
        st.markdown("### Ultimi report")
        data = [{"Data":r["report_date"],"Cliente":r["company_name"],"P.IVA/CF":r["vat_cf"],"Score":f"{r['score']:.0f}/100","Giudizio":r["rating"]} for r in rows]
        st.dataframe(data, use_container_width=True, hide_index=True)


def page_insert() -> None:
    header("Inserisci CR", "Caricamento, estrazione, analisi e generazione report")
    a1, a2 = st.columns([1, 1])
    if a1.button("Nuovo caricamento / pulisci", use_container_width=True):
        for k in ("source_hash", "source_name", "source_bytes", "extracted_text", "client", "analysis", "preview_pdf"):
            st.session_state.pop(k, None)
        st.rerun()
    a2.button(
        "Vai all'elenco report",
        use_container_width=True,
        on_click=request_navigation,
        args=("Elenco Report",),
    )
    up = st.file_uploader("Inserisci Centrale Rischi, documento o allegato (qualsiasi formato)", type=None, accept_multiple_files=False, help="Sono ammessi tutti i formati. PDF, Word, Excel, CSV, testo e immagini possono essere anche analizzati automaticamente; gli altri vengono comunque archiviati.")
    if up is not None:
        b = up.getvalue(); h = sha256(b)
        st.caption(f"{up.name} - {len(b)/1024:.1f} KB - SHA256 {h[:12]}…")
        if st.session_state.get("source_hash") != h:
            if st.button("Analizza documento", type="primary", use_container_width=True):
                with st.spinner("Estrazione dati e calcolo indicatori..."):
                    text, warnings = DocumentExtractor().extract(up.name, b)
                    client, analysis = CRParser().parse(text, warnings)
                    st.session_state.update(source_hash=h, source_name=up.name, source_bytes=b, extracted_text=text, client=asdict(client), analysis=asdict(analysis))
                    st.success("Analisi completata. Verificare i dati estratti.")
                    st.rerun()
    if "analysis" not in st.session_state:
        st.info("Caricare un documento e premere Analizza documento.")
        return
    a_dict = st.session_state["analysis"]
    a = CRAnalysis(**a_dict)
    for w in a.warnings: st.warning(w)
    tabs = st.tabs(["Anagrafica", "Sintesi", "Dati estratti", "Anteprima PDF"])
    with tabs[0]:
        d = st.session_state["client"]
        with st.form("anag"):
            c1,c2 = st.columns(2)
            company = c1.text_input("Denominazione sociale", d.get("company_name",""))
            vat = c2.text_input("Partita IVA / Codice fiscale (facoltativo)", d.get("vat_cf",""), help="Il salvataggio e consentito anche se il campo resta vuoto.")
            tax = c1.text_input("Codice fiscale", d.get("tax_code",""))
            period = c2.text_input("Periodo censito", d.get("period",""))
            address = st.text_input("Via / indirizzo", d.get("address",""))
            c3,c4,c5 = st.columns([2,1,1])
            city = c3.text_input("Comune", d.get("city","")); prov = c4.text_input("Provincia", d.get("province",""), max_chars=2); cciaa = c5.text_input("CCIAA/REA", d.get("cciaa",""))
            lei = st.text_input("Codice LEI", d.get("lei",""))
            apply = st.form_submit_button("Applica dati", use_container_width=True)
        if apply:
            st.session_state["client"] = asdict(ClientData(company.strip(), vat.strip(), tax.strip(), address.strip(), city.strip(), prov.strip().upper(), cciaa.strip(), lei.strip(), period.strip()))
            st.session_state.pop("preview_pdf", None)
            st.success("Anagrafica aggiornata.")
            st.rerun()
    with tabs[1]:
        c1,c2,c3,c4 = st.columns(4)
        c1.metric("Score", f"{a.score:.0f}/100"); c2.metric("Giudizio", a.rating); c3.metric("Affidamenti", money(a.metrics['accordato'])); c4.metric("Utilizzo", pct(a.metrics['utilizzo_pct']))
        st.write(a.narrative)
        st.markdown("#### Alert")
        for x in a.anomalies: st.write("-", x)
        st.markdown("#### Azioni consigliate")
        for x in a.recommendations: st.write("-", x)
    with tabs[2]:
        st.json({"metriche":a.metrics,"banche":a.banks,"categorie":a.categories,"periodi":a.monthly})
        with st.expander("Testo estratto"):
            st.text_area("", st.session_state.get("extracted_text","")[:150000], height=500)
    with tabs[3]:
        if st.button("Genera / aggiorna anteprima", type="primary", use_container_width=True):
            c = ClientData(**st.session_state["client"])
            if not c.company_name.strip():
                c.company_name = "Cliente senza denominazione"
                st.session_state["client"] = asdict(c)
            with st.spinner("Generazione PDF..."):
                st.session_state["preview_pdf"] = build_report(c,a)
        if st.session_state.get("preview_pdf"):
            pdf = st.session_state["preview_pdf"]
            c = ClientData(**st.session_state["client"])
            st.download_button("Scarica anteprima PDF", pdf, file_name=f"{safe_name(c.company_name)}_ANALISI_CR_{it_date()}.pdf", mime="application/pdf", use_container_width=True)
            if st.button("Salva report nell'archivio", use_container_width=True):
                try:
                    created, path = save_report(c,a,st.session_state["source_name"],st.session_state["source_bytes"],pdf)
                    st.success(("Nuovo cliente creato. " if created else "Cliente già esistente: nuovo report aggiunto. ") + f"Salvato in {path}")
                except Exception as e:
                    st.error(str(e))
            pdf_preview(pdf)


def page_reports() -> None:
    header("Elenco Report", "Ricerca, anteprima, download ed eliminazione")
    q = st.text_input("Cerca cliente, P.IVA/CF o giudizio")
    rows = report_rows(q)
    if not rows:
        st.info("Nessun report trovato.")
        return
    labels = {f"{r['company_name']} | {r['vat_cf'] or 'P.IVA/CF non indicata'} | {r['report_date']} | {r['score']:.0f}/100 {r['rating']}": r for r in rows}
    choice = st.selectbox("Seleziona report", list(labels))
    r = labels[choice]
    c1,c2,c3 = st.columns(3); c1.metric("Score",f"{r['score']:.0f}/100"); c2.metric("Giudizio",r['rating']); c3.metric("Data",r['report_date'])
    p = Path(r["report_path"])
    if p.exists():
        data = p.read_bytes()
        st.download_button("Scarica PDF",data,file_name=p.name,mime="application/pdf",use_container_width=True)
        pdf_preview(data)
    else:
        st.error("File PDF non presente nel percorso archivio.")
    st.markdown("---")
    confirm = st.checkbox("Confermo l'eliminazione del report e dei file associati")
    if st.button("Elimina report", disabled=not confirm):
        if delete_report(int(r["id"])):
            st.success("Report eliminato."); st.rerun()


def main() -> None:
    st.set_page_config(page_title=f"{APP_NAME} | FinancePlus.Tech", page_icon="📊", layout="wide", initial_sidebar_state="expanded")
    init_db(); css()
    apply_pending_navigation()
    with st.sidebar:
        st.markdown('<div class="fp-brand">FINANCE<span>PLUS</span><br><small>CR VALUE 2.3</small></div>', unsafe_allow_html=True)
        st.caption(SITE)
        st.markdown("---")
        if st.session_state.get("nav_page") not in {"Dashboard", "Inserisci CR", "Elenco Report"}:
            st.session_state["nav_page"] = "Dashboard"
        page = st.radio(
            "Navigazione",
            ["Dashboard", "Inserisci CR", "Elenco Report"],
            key="nav_page",
            label_visibility="collapsed",
        )
        st.markdown("---")
        st.caption(f"Versione {APP_VERSION}")
        st.caption("Score consulenziale interno - verificare sempre la fonte originale.")
    if page == "Dashboard": page_dashboard()
    elif page == "Inserisci CR": page_insert()
    else: page_reports()


if __name__ == "__main__":
    main()
